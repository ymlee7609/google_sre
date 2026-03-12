#!/usr/bin/env python3
"""SRE 교차 학습 진도표 PPTX 빌더 - Glassmorphism 디자인.

output/SRE_Study_Guide.md 의 교차 학습 진도표를
Glassmorphism 스타일의 PPTX로 변환합니다.
"""

import logging
from pathlib import Path

from lxml.etree import SubElement
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

BASE_DIR = Path(__file__).resolve().parent.parent
OUTPUT_DIR = BASE_DIR / "output"


# --- Glassmorphism 디자인 상수 ---
class GlassTheme:
    """Glassmorphism 디자인 사양."""

    # 배경
    BG_DEEP = RGBColor(0x0F, 0x0F, 0x2D)

    # 글로우 블롭 색상
    GLOW_VIOLET = RGBColor(0x6B, 0x21, 0xA8)
    GLOW_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
    GLOW_CYAN = RGBColor(0x06, 0x7A, 0x8A)

    # 텍스트
    TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_SOFT = RGBColor(0xE0, 0xE0, 0xF0)
    TEXT_DIM = RGBColor(0xA0, 0xA0, 0xC0)

    # 악센트
    ACCENT_CYAN = RGBColor(0x67, 0xE8, 0xF9)
    ACCENT_VIOLET = RGBColor(0xA7, 0x8B, 0xFA)

    # 글래스 카드
    GLASS_FILL = RGBColor(0xFF, 0xFF, 0xFF)  # 15-20% opacity
    GLASS_BORDER = RGBColor(0xFF, 0xFF, 0xFF)  # 25% opacity

    # 폰트
    FONT_TITLE = "Calibri Light"
    FONT_BODY = "Segoe UI"
    FONT_MONO = "Consolas"

    # 슬라이드 크기 (16:9)
    SLIDE_WIDTH = Emu(12192000)
    SLIDE_HEIGHT = Emu(6858000)

    # 레이아웃
    MARGIN_LEFT = Inches(0.8)
    CONTENT_WIDTH = Inches(11.733)


# --- 슬라이드 빌더 ---
class GlassSlideBuilder:
    """Glassmorphism 스타일 슬라이드 생성."""

    def __init__(self, prs: Presentation) -> None:
        self.prs = prs
        self.blank_layout = prs.slide_layouts[6]

    # --- 유틸리티 ---

    @staticmethod
    def _set_transparency(shape, pct: int) -> None:
        """도형 채우기 투명도 설정."""
        alpha_val = str((100 - pct) * 1000)
        for sf in shape._element.iter(qn("a:solidFill")):
            for color_el in sf:
                alpha_el = SubElement(color_el, qn("a:alpha"))
                alpha_el.set("val", alpha_val)
            break

    @staticmethod
    def _set_line_transparency(shape, pct: int) -> None:
        """도형 외곽선 투명도 설정."""
        alpha_val = str((100 - pct) * 1000)
        ln = shape._element.find(qn("a:ln"), shape._element.nsmap)
        if ln is None:
            ln = shape.line._ln
        for sf in ln.iter(qn("a:solidFill")):
            for color_el in sf:
                alpha_el = SubElement(color_el, qn("a:alpha"))
                alpha_el.set("val", alpha_val)
            break

    def _new_slide(self):
        """새 슬라이드 생성 + 배경 + 글로우."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        # 배경색
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = GlassTheme.BG_DEEP
        # 글로우 블롭
        self._add_glow(slide, Inches(-1.5), Inches(-1.0),
                       Inches(6), Inches(5), GlassTheme.GLOW_VIOLET, 75)
        self._add_glow(slide, Inches(8.5), Inches(4.0),
                       Inches(6), Inches(4.5), GlassTheme.GLOW_BLUE, 80)
        self._add_glow(slide, Inches(4.0), Inches(5.5),
                       Inches(4), Inches(3), GlassTheme.GLOW_CYAN, 85)
        return slide

    def _add_glow(self, slide, left, top, w, h, color, transparency) -> None:
        """배경 글로우 블롭."""
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        self._set_transparency(shape, transparency)

    def _add_glass_card(self, slide, left, top, width, height):
        """글래스 카드 (반투명 사각형)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = GlassTheme.GLASS_FILL
        self._set_transparency(shape, 85)
        shape.line.color.rgb = GlassTheme.GLASS_BORDER
        shape.line.width = Pt(0.75)
        self._set_line_transparency(shape, 75)
        return shape

    def _add_accent_line(self, slide, left, top, width) -> None:
        """악센트 수평선."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Pt(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = GlassTheme.ACCENT_CYAN
        shape.line.fill.background()
        self._set_transparency(shape, 40)

    def _add_text(self, slide, left, top, width, height, text,
                  font_name=None, font_size=None, color=None,
                  bold=False, alignment=None) -> None:
        """텍스트 박스 추가."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.name = font_name or GlassTheme.FONT_BODY
        p.font.size = font_size or Pt(14)
        p.font.color.rgb = color or GlassTheme.TEXT_SOFT
        p.font.bold = bold
        if alignment:
            p.alignment = alignment
        return tb

    def _add_bullet_list(self, slide, left, top, width, height,
                         items, font_size=None, color=None) -> None:
        """불릿 리스트 추가."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        fs = font_size or Pt(13)
        c = color or GlassTheme.TEXT_SOFT

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"\u2022  {item}"
            run.font.name = GlassTheme.FONT_BODY
            run.font.size = fs
            run.font.color.rgb = c

    def _add_footer(self, slide, text: str) -> None:
        """하단 캡션."""
        self._add_accent_line(slide, Inches(0.5), Inches(6.85), Inches(12.333))
        self._add_text(
            slide, Inches(0.5), Inches(6.95), Inches(10), Inches(0.3),
            text, GlassTheme.FONT_MONO, Pt(8), GlassTheme.TEXT_DIM,
        )

    # --- 슬라이드 유형 ---

    def add_cover(self) -> None:
        """커버 슬라이드."""
        slide = self._new_slide()

        # 중앙 글래스 카드
        self._add_glass_card(
            slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5),
        )

        # 부제
        self._add_text(
            slide, Inches(2.0), Inches(2.0), Inches(9.0), Inches(0.5),
            "Google SRE 3\uad8c \uad50\ucc28 \ud559\uc2b5",
            GlassTheme.FONT_BODY, Pt(16), GlassTheme.ACCENT_CYAN,
        )

        # 제목
        self._add_text(
            slide, Inches(2.0), Inches(2.7), Inches(9.0), Inches(1.2),
            "SRE Study Guide",
            GlassTheme.FONT_TITLE, Pt(48), GlassTheme.TEXT_WHITE, bold=True,
        )

        self._add_accent_line(slide, Inches(2.0), Inches(4.0), Inches(3.0))

        # 설명
        self._add_text(
            slide, Inches(2.0), Inches(4.3), Inches(9.0), Inches(1.0),
            "12\uc8fc \ucee4\ub9ac\ud058\ub7fc  |  34 + 21 + 21 Chapters  |  \uad50\ucc28 \ucc38\uc870 \ud14c\uc774\ube14",
            GlassTheme.FONT_BODY, Pt(14), GlassTheme.TEXT_DIM,
        )

        self._add_footer(slide, "Site Reliability Engineering  \u2014  Cross-Study Guide")

    def add_overview(self) -> None:
        """3\uad8c \uc5ed\ud560 \uac1c\uc694 \uc2ac\ub77c\uc774\ub4dc."""
        slide = self._new_slide()

        self._add_text(
            slide, GlassTheme.MARGIN_LEFT, Inches(0.5), Inches(10), Inches(0.6),
            "3\uad8c\uc758 \uc5ed\ud560",
            GlassTheme.FONT_TITLE, Pt(36), GlassTheme.TEXT_WHITE, bold=True,
        )
        self._add_accent_line(slide, GlassTheme.MARGIN_LEFT, Inches(1.15), Inches(2.5))

        books = [
            ("SRE Book", "\uc8fc\uad50\uc7ac (\uc774\ub860)", "Google SRE\uc758 \uc6d0\uce59, \uc2e4\ubb34, \uc870\uc9c1 \uc6b4\uc601",
             "34 Chapters + 6 Appendices", GlassTheme.ACCENT_CYAN),
            ("Workbook", "\uc2e4\uc2b5 \ubcf4\ucda9", "SRE Book\uc758 \uac1c\ub150\uc744 \uc2e4\ubb34\uc5d0 \uc801\uc6a9\ud558\ub294 \ubc29\ubc95",
             "21 Chapters + 3 Appendices", GlassTheme.ACCENT_VIOLET),
            ("BSRS", "\ubcf4\uc548/\uc124\uacc4 \ubcf4\ucda9", "\ubcf4\uc548\uacfc \uc2e0\ub8b0\uc131\uc744 \ud1b5\ud569\ud55c \uc2dc\uc2a4\ud15c \uc124\uacc4",
             "21 Chapters + 1 Appendix", RGBColor(0x6E, 0xE7, 0xB7)),
        ]

        for i, (name, role, desc, stats, accent) in enumerate(books):
            card_top = Inches(1.6 + i * 1.7)
            self._add_glass_card(
                slide, Inches(1.0), card_top, Inches(11.333), Inches(1.4),
            )
            # 도서명
            self._add_text(
                slide, Inches(1.4), card_top + Inches(0.15),
                Inches(3.0), Inches(0.5),
                name, GlassTheme.FONT_TITLE, Pt(24), accent, bold=True,
            )
            # 역할
            self._add_text(
                slide, Inches(4.5), card_top + Inches(0.2),
                Inches(2.5), Inches(0.4),
                role, GlassTheme.FONT_BODY, Pt(14), GlassTheme.TEXT_DIM,
            )
            # 설명
            self._add_text(
                slide, Inches(1.4), card_top + Inches(0.7),
                Inches(7.0), Inches(0.5),
                desc, GlassTheme.FONT_BODY, Pt(13), GlassTheme.TEXT_SOFT,
            )
            # 통계
            self._add_text(
                slide, Inches(8.5), card_top + Inches(0.7),
                Inches(3.5), Inches(0.4),
                stats, GlassTheme.FONT_MONO, Pt(11), GlassTheme.TEXT_DIM,
            )

        self._add_footer(slide, "SRE Study Guide  \u2014  Overview")

    def add_method(self) -> None:
        """\ud559\uc2b5 \ubc29\ubc95 \uc2ac\ub77c\uc774\ub4dc."""
        slide = self._new_slide()

        self._add_text(
            slide, GlassTheme.MARGIN_LEFT, Inches(0.5), Inches(10), Inches(0.6),
            "\ud559\uc2b5 \ubc29\ubc95",
            GlassTheme.FONT_TITLE, Pt(36), GlassTheme.TEXT_WHITE, bold=True,
        )
        self._add_accent_line(slide, GlassTheme.MARGIN_LEFT, Inches(1.15), Inches(2.5))

        steps = [
            ("01", "SRE Book \ud559\uc2b5", "SRE Book\uc744 \uc8fc\uad50\uc7ac\ub85c \uc21c\uc11c\ub300\ub85c \ud559\uc2b5"),
            ("02", "\uad50\ucc28 \ucc38\uc870 \ud655\uc778", "\uac01 \ucc55\ud130 \ud559\uc2b5 \ud6c4 \uad00\ub828 Workbook/BSRS \ucc55\ud130 \ud655\uc778"),
            ("03", "Workbook \ubcf4\ucda9", "Workbook\uc73c\ub85c \uc2e4\ubb34 \uc801\uc6a9 \ubc29\ubc95\uc744 \ubcf4\ucda9 \ud559\uc2b5"),
            ("04", "BSRS \ubcf4\ucda9", "BSRS\ub85c \ubcf4\uc548/\uc2e0\ub8b0\uc131 \uad00\uc810\uc5d0\uc11c \uc124\uacc4 \ubcf4\ucda9"),
        ]

        for i, (num, title, desc) in enumerate(steps):
            card_top = Inches(1.5 + i * 1.25)
            self._add_glass_card(
                slide, Inches(1.0), card_top, Inches(11.333), Inches(1.0),
            )
            # 번호
            self._add_text(
                slide, Inches(1.3), card_top + Inches(0.15),
                Inches(0.8), Inches(0.7),
                num, GlassTheme.FONT_TITLE, Pt(28), GlassTheme.ACCENT_CYAN, bold=True,
            )
            # 제목
            self._add_text(
                slide, Inches(2.3), card_top + Inches(0.1),
                Inches(4.0), Inches(0.4),
                title, GlassTheme.FONT_TITLE, Pt(20), GlassTheme.TEXT_WHITE, bold=True,
            )
            # 설명
            self._add_text(
                slide, Inches(2.3), card_top + Inches(0.55),
                Inches(9.0), Inches(0.4),
                desc, GlassTheme.FONT_BODY, Pt(13), GlassTheme.TEXT_SOFT,
            )

        self._add_footer(slide, "SRE Study Guide  \u2014  Study Method")

    def add_curriculum_overview(self) -> None:
        """12\uc8fc \ucee4\ub9ac\ud058\ub7fc \uc804\uccb4 \uac1c\uc694."""
        slide = self._new_slide()

        self._add_text(
            slide, GlassTheme.MARGIN_LEFT, Inches(0.4), Inches(10), Inches(0.6),
            "12\uc8fc \ud559\uc2b5 \ucee4\ub9ac\ud058\ub7fc",
            GlassTheme.FONT_TITLE, Pt(32), GlassTheme.TEXT_WHITE, bold=True,
        )
        self._add_accent_line(slide, GlassTheme.MARGIN_LEFT, Inches(1.0), Inches(2.5))

        weeks = [
            ("W1", "Part I: \uc11c\ub860", "2Ch", GlassTheme.ACCENT_CYAN),
            ("W2-3", "Part II: \uc6d0\uce59", "7Ch", GlassTheme.ACCENT_CYAN),
            ("W4-8", "Part III: \uc2e4\ubb34", "18Ch", GlassTheme.ACCENT_VIOLET),
            ("W9-10", "Part IV: \uad00\ub9ac", "5Ch", GlassTheme.ACCENT_VIOLET),
            ("W11", "Part V: \uacb0\ub860/\ubd80\ub85d", "2Ch+App", RGBColor(0x6E, 0xE7, 0xB7)),
            ("W12", "\ubcf4\ucda9 \ud559\uc2b5", "9Ch", RGBColor(0x6E, 0xE7, 0xB7)),
        ]

        for i, (week, part, count, accent) in enumerate(weeks):
            row = i // 3
            col = i % 3
            left = Inches(0.6 + col * 4.1)
            top = Inches(1.3 + row * 2.7)

            self._add_glass_card(slide, left, top, Inches(3.8), Inches(2.3))

            # 주차
            self._add_text(
                slide, left + Inches(0.3), top + Inches(0.2),
                Inches(1.5), Inches(0.5),
                week, GlassTheme.FONT_TITLE, Pt(28), accent, bold=True,
            )
            # 챕터 수
            self._add_text(
                slide, left + Inches(2.5), top + Inches(0.25),
                Inches(1.0), Inches(0.4),
                count, GlassTheme.FONT_MONO, Pt(14), GlassTheme.TEXT_DIM,
            )
            # Part 이름
            self._add_text(
                slide, left + Inches(0.3), top + Inches(0.8),
                Inches(3.2), Inches(0.5),
                part, GlassTheme.FONT_BODY, Pt(16), GlassTheme.TEXT_WHITE,
            )

        self._add_footer(slide, "SRE Study Guide  \u2014  12-Week Curriculum")

    def add_week_slide(self, week_num: str, title: str,
                       rows: list[tuple[str, str, str, str]]) -> None:
        """개별 주차 진도표 슬라이드."""
        slide = self._new_slide()

        self._add_text(
            slide, GlassTheme.MARGIN_LEFT, Inches(0.4), Inches(10), Inches(0.6),
            f"Week {week_num}",
            GlassTheme.FONT_TITLE, Pt(32), GlassTheme.ACCENT_CYAN, bold=True,
        )
        self._add_text(
            slide, Inches(3.5), Inches(0.45), Inches(8), Inches(0.5),
            title, GlassTheme.FONT_BODY, Pt(18), GlassTheme.TEXT_DIM,
        )
        self._add_accent_line(slide, GlassTheme.MARGIN_LEFT, Inches(1.0), Inches(11.5))

        # 테이블 영역 글래스 카드
        card_top = Inches(1.2)
        card_height = min(Inches(0.65) * len(rows) + Inches(0.7), Inches(5.5))
        self._add_glass_card(
            slide, Inches(0.5), card_top, Inches(12.333), card_height,
        )

        # 헤더
        headers = [("\uc77c\ucc28", 0.8), ("SRE Book", 3.5),
                    ("\ubcf4\ucda9 \uc790\ub8cc", 3.5), ("\ud559\uc2b5 \ud3ec\uc778\ud2b8", 3.5)]
        x = Inches(0.7)
        for hdr, w in headers:
            self._add_text(
                slide, x, card_top + Inches(0.1), Inches(w), Inches(0.4),
                hdr, GlassTheme.FONT_BODY, Pt(11), GlassTheme.ACCENT_VIOLET, bold=True,
            )
            x += Inches(w)

        # 데이터 행
        for i, (day, sre, supp, point) in enumerate(rows):
            y = card_top + Inches(0.55 + i * 0.6)
            x = Inches(0.7)
            cols = [(day, 0.8), (sre, 3.5), (supp, 3.5), (point, 3.5)]
            for val, w in cols:
                color = GlassTheme.TEXT_WHITE if w == 0.8 else GlassTheme.TEXT_SOFT
                fs = Pt(11) if w == 0.8 else Pt(11)
                self._add_text(slide, x, y, Inches(w), Inches(0.5),
                               val, GlassTheme.FONT_BODY, fs, color)
                x += Inches(w)

        self._add_footer(slide, f"SRE Study Guide  \u2014  Week {week_num}")

    def add_cross_ref_slide(self, part_name: str, part_num: str,
                            rows: list[tuple[str, str, str]]) -> None:
        """교차 참조 테이블 슬라이드."""
        slide = self._new_slide()

        self._add_text(
            slide, GlassTheme.MARGIN_LEFT, Inches(0.35), Inches(6), Inches(0.5),
            f"\uad50\ucc28 \ucc38\uc870  \u2014  {part_name}",
            GlassTheme.FONT_TITLE, Pt(28), GlassTheme.TEXT_WHITE, bold=True,
        )
        self._add_accent_line(slide, GlassTheme.MARGIN_LEFT, Inches(0.9), Inches(2.5))

        card_top = Inches(1.05)
        row_h = 0.5
        card_height = min(row_h * len(rows) + 0.65, 5.6)
        self._add_glass_card(
            slide, Inches(0.4), card_top,
            Inches(12.533), Inches(card_height),
        )

        # 헤더
        hdr_cols = [("SRE Book", Inches(0.6), 4.0),
                    ("Workbook", Inches(4.6), 3.8),
                    ("BSRS", Inches(8.5), 3.8)]
        for hdr, x, w in hdr_cols:
            self._add_text(
                slide, x, card_top + Inches(0.08), Inches(w), Inches(0.35),
                hdr, GlassTheme.FONT_BODY, Pt(11), GlassTheme.ACCENT_CYAN, bold=True,
            )

        # 데이터 - 폰트 크기 조정
        fs = Pt(10) if len(rows) > 8 else Pt(11)
        rh = 0.45 if len(rows) > 8 else row_h

        for i, (sre, wb, bsrs) in enumerate(rows):
            y = card_top + Inches(0.45 + i * rh)
            self._add_text(slide, Inches(0.6), y, Inches(4.0), Inches(0.4),
                           sre, GlassTheme.FONT_BODY, fs, GlassTheme.TEXT_WHITE)
            self._add_text(slide, Inches(4.6), y, Inches(3.8), Inches(0.4),
                           wb, GlassTheme.FONT_BODY, fs, GlassTheme.TEXT_SOFT)
            self._add_text(slide, Inches(8.5), y, Inches(3.8), Inches(0.4),
                           bsrs, GlassTheme.FONT_BODY, fs, GlassTheme.TEXT_SOFT)

        self._add_footer(slide, f"SRE Study Guide  \u2014  Cross Reference {part_num}")

    def add_supplement_slide(self, title: str,
                             items: list[tuple[str, str]]) -> None:
        """보충 학습 슬라이드."""
        slide = self._new_slide()

        self._add_text(
            slide, GlassTheme.MARGIN_LEFT, Inches(0.4), Inches(10), Inches(0.6),
            title,
            GlassTheme.FONT_TITLE, Pt(32), GlassTheme.TEXT_WHITE, bold=True,
        )
        self._add_accent_line(slide, GlassTheme.MARGIN_LEFT, Inches(1.05), Inches(2.5))

        for i, (ch, desc) in enumerate(items):
            top = Inches(1.3 + i * 0.85)
            self._add_glass_card(slide, Inches(0.6), top, Inches(12.133), Inches(0.7))
            self._add_text(
                slide, Inches(0.9), top + Inches(0.1),
                Inches(6.0), Inches(0.4),
                ch, GlassTheme.FONT_BODY, Pt(13), GlassTheme.ACCENT_VIOLET,
            )
            self._add_text(
                slide, Inches(7.0), top + Inches(0.1),
                Inches(5.5), Inches(0.4),
                desc, GlassTheme.FONT_BODY, Pt(12), GlassTheme.TEXT_SOFT,
            )

        self._add_footer(slide, "SRE Study Guide  \u2014  Supplementary")

    def add_study_paths(self) -> None:
        """주제별 집중 학습 경로."""
        slide = self._new_slide()

        self._add_text(
            slide, GlassTheme.MARGIN_LEFT, Inches(0.4), Inches(10), Inches(0.6),
            "\uc8fc\uc81c\ubcc4 \uc9d1\uc911 \ud559\uc2b5 \uacbd\ub85c",
            GlassTheme.FONT_TITLE, Pt(32), GlassTheme.TEXT_WHITE, bold=True,
        )
        self._add_accent_line(slide, GlassTheme.MARGIN_LEFT, Inches(1.05), Inches(2.5))

        paths = [
            ("SLO / \ubaa8\ub2c8\ud130\ub9c1",
             "SRE Ch03-06 \u2192 Workbook Ch02-05 \u2192 BSRS Ch04, Ch15",
             GlassTheme.ACCENT_CYAN),
            ("\uc778\uc2dc\ub358\ud2b8 \ub300\uc751",
             "SRE Ch13-15 \u2192 Workbook Ch09-10 \u2192 BSRS Ch17-18",
             GlassTheme.ACCENT_VIOLET),
            ("\ubcf4\uc548 \uc911\uc2ec",
             "BSRS Ch01-05 \u2192 SRE Ch17-18 \u2192 BSRS Ch12-14",
             RGBColor(0xFF, 0x6B, 0x6B)),
            ("\ub300\uaddc\ubaa8 \uc2dc\uc2a4\ud15c \uc124\uacc4",
             "SRE Ch19-24 \u2192 Workbook Ch11-13 \u2192 BSRS Ch08-09",
             RGBColor(0x6E, 0xE7, 0xB7)),
            ("\uc870\uc9c1 / \ubb38\ud654",
             "SRE Ch28-32 \u2192 Workbook Ch18-21 \u2192 BSRS Ch20-21",
             RGBColor(0xFF, 0xD9, 0x3D)),
        ]

        for i, (title, path, accent) in enumerate(paths):
            top = Inches(1.3 + i * 1.05)
            self._add_glass_card(slide, Inches(0.6), top, Inches(12.133), Inches(0.85))
            self._add_text(
                slide, Inches(1.0), top + Inches(0.08),
                Inches(4.0), Inches(0.35),
                title, GlassTheme.FONT_TITLE, Pt(18), accent, bold=True,
            )
            self._add_text(
                slide, Inches(1.0), top + Inches(0.45),
                Inches(11.0), Inches(0.35),
                path, GlassTheme.FONT_MONO, Pt(12), GlassTheme.TEXT_SOFT,
            )

        self._add_footer(slide, "SRE Study Guide  \u2014  Study Paths")

    def add_closing(self) -> None:
        """클로징 슬라이드."""
        slide = self._new_slide()

        self._add_glass_card(
            slide, Inches(2.5), Inches(2.0), Inches(8.333), Inches(3.5),
        )

        self._add_text(
            slide, Inches(3.0), Inches(2.5), Inches(7.333), Inches(1.0),
            "Happy SRE Learning!",
            GlassTheme.FONT_TITLE, Pt(44), GlassTheme.TEXT_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        self._add_accent_line(slide, Inches(5.0), Inches(3.7), Inches(3.333))

        self._add_text(
            slide, Inches(3.0), Inches(4.0), Inches(7.333), Inches(0.8),
            "SRE Book \u00b7 Workbook \u00b7 BSRS\n12\uc8fc \uad50\ucc28 \ud559\uc2b5 \ucee4\ub9ac\ud058\ub7fc",
            GlassTheme.FONT_BODY, Pt(16), GlassTheme.TEXT_DIM,
            alignment=PP_ALIGN.CENTER,
        )

        self._add_footer(slide, "Site Reliability Engineering  \u2014  Cross-Study Guide")


# --- 콘텐츠 데이터 ---

WEEKS = [
    ("1", "Part I: \uc11c\ub860", [
        ("1\uc77c", "Ch01 Introduction", "Workbook Ch01, BSRS Ch01",
         "SRE \uc815\uc758, DevOps \uad00\uacc4, \ubcf4\uc548-\uc2e0\ub8b0\uc131 \uad50\ucc28\uc810"),
        ("2\uc77c", "Ch02 Production Environment", "\u2014",
         "Google \ud504\ub85c\ub355\uc158 \ud658\uacbd \uc774\ud574"),
    ]),
    ("2", "Part II: \uc6d0\uce59 (\uc804\ubc18)", [
        ("1\uc77c", "Ch03 Embracing Risk", "WB Ch02-03, BSRS Ch04",
         "\ub9ac\uc2a4\ud06c \uc218\uc6a9, SLO \uad6c\ud604, \uc124\uacc4 \ud2b8\ub808\uc774\ub4dc\uc624\ud504"),
        ("2\uc77c", "Ch04 Service Level Objectives", "WB Ch02-03, Ch05",
         "SLO \uc815\uc758, SLO \uae30\ubc18 \uc54c\ub9bc"),
        ("3\uc77c", "Ch05 Eliminating Toil", "WB Ch06",
         "\ud1a0\uc77c \uc815\uc758, \uc81c\uac70 \uc804\ub7b5"),
        ("4\uc77c", "Ch06 Monitoring", "WB Ch04-05, BSRS Ch15",
         "\ubaa8\ub2c8\ud130\ub9c1 \uc6d0\uce59, SLO \uc54c\ub9bc"),
    ]),
    ("3", "Part II: \uc6d0\uce59 (\ud6c4\ubc18)", [
        ("1\uc77c", "Ch07 Automation", "BSRS Ch07",
         "\uc790\ub3d9\ud654 \ubc1c\uc804, \ubcc0\ud654 \ub300\uc751 \uc124\uacc4"),
        ("2\uc77c", "Ch08 Release Engineering", "WB Ch16, BSRS Ch14",
         "\ub9b4\ub9ac\uc2a4 \uc5d4\uc9c0\ub2c8\uc5b4\ub9c1, \uce74\ub098\ub9ac \ubc30\ud3ec"),
        ("3\uc77c", "Ch09 Simplicity", "WB Ch07, BSRS Ch06",
         "\ub2e8\uc21c\uc131 \uc6d0\uce59, \uc774\ud574 \uac00\ub2a5\ud55c \uc124\uacc4"),
    ]),
    ("4", "Part III: \uc2e4\ubb34 (\uc54c\ub9bc & \uc628\ucf5c)", [
        ("1\uc77c", "Ch10 Practical Alerting", "WB Ch04-05",
         "\uc2dc\uacc4\uc5f4 \uae30\ubc18 \uc54c\ub9bc \uc2e4\ubb34"),
        ("2\uc77c", "Ch11 Being On-Call", "WB Ch08",
         "\uc628\ucf5c \uc6b4\uc601 \uc6d0\uce59\uacfc \uc2e4\ubb34"),
        ("3\uc77c", "Ch12 Troubleshooting", "BSRS Ch15",
         "\ud6a8\uacfc\uc801\uc778 \ubb38\uc81c \ud574\uacb0"),
    ]),
    ("5", "Part III: \uc778\uc2dc\ub358\ud2b8 \ub300\uc751", [
        ("1\uc77c", "Ch13 Emergency Response", "WB Ch09, BSRS Ch17",
         "\uae34\uae09 \ub300\uc751, \uc704\uae30 \uad00\ub9ac"),
        ("2\uc77c", "Ch14 Managing Incidents", "WB Ch09, BSRS Ch17",
         "\uc778\uc2dc\ub358\ud2b8 \uad00\ub9ac \ud504\ub85c\uc138\uc2a4"),
        ("3\uc77c", "Ch15 Postmortem Culture", "WB Ch10, BSRS Ch18",
         "\ud3ec\uc2a4\ud2b8\ubaa8\ud15c \ubb38\ud654, \ubcf5\uad6c \ud6c4\uc18d \uc870\uce58"),
        ("4\uc77c", "Ch16 Tracking Outages", "\u2014",
         "\uc7a5\uc560 \ucd94\uc801 \uc2dc\uc2a4\ud15c"),
    ]),
    ("6", "Part III: \ud14c\uc2a4\ud2b8 & \uac1c\ubc1c", [
        ("1\uc77c", "Ch17 Testing for Reliability", "BSRS Ch13",
         "\uc2e0\ub8b0\uc131 \ud14c\uc2a4\ud2b8, \ucf54\ub4dc \ud14c\uc2a4\ud2b8"),
        ("2\uc77c", "Ch18 Software Engineering in SRE", "BSRS Ch12",
         "SRE\uc758 SW \uc5d4\uc9c0\ub2c8\uc5b4\ub9c1, \ucf54\ub4dc \uc791\uc131"),
        ("3\uc77c", "Ch19 Load Balancing (FE)", "WB Ch11",
         "\ud504\ub860\ud2b8\uc5d4\ub4dc \ub85c\ub4dc \ubc38\ub7f0\uc2f1"),
    ]),
    ("7", "Part III: \ubd80\ud558 & \uc7a5\uc560", [
        ("1\uc77c", "Ch20 Load Balancing (DC)", "WB Ch11",
         "\ub370\uc774\ud130\uc13c\ud130 \ub85c\ub4dc \ubc38\ub7f0\uc2f1"),
        ("2\uc77c", "Ch21 Handling Overload", "WB Ch17, BSRS Ch10",
         "\uacfc\ubd80\ud558 \ucc98\ub9ac, DoS \ub300\uc751"),
        ("3\uc77c", "Ch22 Cascading Failures", "WB Ch17, BSRS Ch08",
         "\uc5f0\uc1c4 \uc7a5\uc560 \ub300\uc751, \ubcf5\uc6d0\ub825 \uc124\uacc4"),
    ]),
    ("8", "Part III: \ubd84\uc0b0 \uc2dc\uc2a4\ud15c & \ub370\uc774\ud130", [
        ("1\uc77c", "Ch23 Managing Critical State", "\u2014",
         "\ubd84\uc0b0 \ud569\uc758, \ud06c\ub9ac\ud2f0\uceec \uc0c1\ud0dc \uad00\ub9ac"),
        ("2\uc77c", "Ch24 Distributed Scheduling", "\u2014",
         "\ubd84\uc0b0 \ud06c\ub860 \uc2a4\ucf00\uc904\ub9c1"),
        ("3\uc77c", "Ch25 Data Pipelines", "WB Ch13",
         "\ub370\uc774\ud130 \ud30c\uc774\ud504\ub77c\uc778 \uc124\uacc4\uc640 \uc6b4\uc601"),
        ("4\uc77c", "Ch26 Data Integrity", "BSRS Ch09",
         "\ub370\uc774\ud130 \ubb34\uacb0\uc131, \ubcf5\uad6c \uc124\uacc4"),
    ]),
    ("9", "Part III \ub9c8\ubb34\ub9ac + Part IV \uc2dc\uc791", [
        ("1\uc77c", "Ch27 Product Launches", "WB Ch16, BSRS Ch14",
         "\ub300\uaddc\ubaa8 \uc81c\ud488 \ucd9c\uc2dc, \uce74\ub098\ub9ac \ubc30\ud3ec"),
        ("2\uc77c", "Ch28 Accelerating SREs", "WB Ch20",
         "SRE \uc628\ubcf4\ub529, \ud300 \uc0dd\uba85\uc8fc\uae30"),
        ("3\uc77c", "Ch29 Dealing with Interrupts", "\u2014",
         "\uc778\ud130\ub7fd\ud2b8 \uad00\ub9ac"),
    ]),
    ("10", "Part IV: \uad00\ub9ac", [
        ("1\uc77c", "Ch30 Operational Overload", "WB Ch18",
         "SRE \ud30c\uacac, \uc6b4\uc601 \uacfc\ubd80\ud558 \ubcf5\uad6c"),
        ("2\uc77c", "Ch31 Communication", "WB Ch19, BSRS Ch20",
         "SRE \ucee4\ubba4\ub2c8\ucf00\uc774\uc158, \uc5ed\ud560\uacfc \ucc45\uc784"),
        ("3\uc77c", "Ch32 SRE Engagement", "WB Ch18, Ch21",
         "SRE \ucc38\uc5ec \ubaa8\ub378 \uc9c4\ud654, \uc870\uc9c1 \ubcc0\ud654"),
    ]),
    ("11", "Part V: \uacb0\ub860 \ubc0f \ubd80\ub85d", [
        ("1\uc77c", "Ch33 Lessons Learned", "BSRS Ch21",
         "\ud0c0 \uc0b0\uc5c5\uc758 \uad50\ud6c8, \ubcf4\uc548/\uc2e0\ub8b0\uc131 \ubb38\ud654"),
        ("2\uc77c", "Ch34 Conclusion", "WB/BSRS Conclusion",
         "3\uad8c \uc885\ud569 \uc815\ub9ac"),
        ("3\uc77c", "Appendix A-F", "WB App A-C, BSRS App A",
         "\ucc38\uc870 \uc790\ub8cc \ud1b5\ud569 \ud559\uc2b5"),
    ]),
    ("12", "\ubcf4\ucda9 \ud559\uc2b5 (\uace0\uc720 \ucc55\ud130)", [
        ("1\uc77c", "BSRS Ch02-03", "\u2014",
         "\uc704\ud611 \ud589\uc704\uc790, Safe Proxies"),
        ("2\uc77c", "BSRS Ch05, Ch11", "\u2014",
         "\ucd5c\uc18c \uad8c\ud55c, \uacf5\uac1c CA"),
        ("3\uc77c", "BSRS Ch16, Ch19", "\u2014",
         "\uc7ac\ud574 \uacc4\ud68d, Chrome \ubcf4\uc548\ud300"),
        ("4\uc77c", "WB Ch12, Ch14-15", "\u2014",
         "NALSD, \uc124\uc815 \uad00\ub9ac \uc124\uacc4"),
    ]),
]

CROSS_REFS = {
    "Part I \u2014 \uc11c\ub860": ("I", [
        ("Ch01 Introduction", "Ch01 How SRE Relates to DevOps",
         "Ch01 Intersection of Security and Reliability"),
        ("Ch02 Production Environment", "\u2014", "\u2014"),
    ]),
    "Part II \u2014 \uc6d0\uce59": ("II", [
        ("Ch03 Embracing Risk", "Ch02-03 SLOs / Case Studies", "Ch04 Design Tradeoffs"),
        ("Ch04 Service Level Objectives", "Ch02-03, Ch05 Alerting on SLOs", "\u2014"),
        ("Ch05 Eliminating Toil", "Ch06 Eliminating Toil", "\u2014"),
        ("Ch06 Monitoring", "Ch04-05 Monitoring / Alerting", "Ch15 Investigating Systems"),
        ("Ch07 Automation", "\u2014", "Ch07 Design for Changing Landscape"),
        ("Ch08 Release Engineering", "Ch16 Canarying Releases", "Ch14 Deploying Code"),
        ("Ch09 Simplicity", "Ch07 Simplicity", "Ch06 Design for Understandability"),
    ]),
    "Part III \u2014 \uc2e4\ubb34 (1/2)": ("III-1", [
        ("Ch10 Practical Alerting", "Ch04-05 Monitoring / Alerting", "\u2014"),
        ("Ch11 Being On-Call", "Ch08 On-Call", "\u2014"),
        ("Ch12 Troubleshooting", "\u2014", "Ch15 Investigating Systems"),
        ("Ch13 Emergency Response", "Ch09 Incident Response", "Ch17 Crisis Management"),
        ("Ch14 Managing Incidents", "Ch09 Incident Response", "Ch17 Crisis Management"),
        ("Ch15 Postmortem Culture", "Ch10 Postmortem Culture", "Ch18 Recovery and Aftermath"),
        ("Ch16 Tracking Outages", "\u2014", "\u2014"),
        ("Ch17 Testing for Reliability", "\u2014", "Ch13 Testing Code"),
        ("Ch18 Software Engineering", "\u2014", "Ch12 Writing Code"),
    ]),
    "Part III \u2014 \uc2e4\ubb34 (2/2)": ("III-2", [
        ("Ch19-20 Load Balancing", "Ch11 Managing Load", "\u2014"),
        ("Ch21 Handling Overload", "Ch17 Overload Recovery", "Ch10 Mitigating DoS"),
        ("Ch22 Cascading Failures", "Ch17 Overload Recovery", "Ch08 Design for Resilience"),
        ("Ch23 Managing Critical State", "\u2014", "\u2014"),
        ("Ch24 Distributed Scheduling", "\u2014", "\u2014"),
        ("Ch25 Data Pipelines", "Ch13 Data Processing", "\u2014"),
        ("Ch26 Data Integrity", "\u2014", "Ch09 Design for Recovery"),
        ("Ch27 Product Launches", "Ch16 Canarying Releases", "Ch14 Deploying Code"),
    ]),
    "Part IV \u2014 \uad00\ub9ac": ("IV", [
        ("Ch28 Accelerating SREs", "Ch20 SRE Team Lifecycles", "\u2014"),
        ("Ch29 Dealing with Interrupts", "\u2014", "\u2014"),
        ("Ch30 Operational Overload", "Ch18 SRE Engagement Model", "\u2014"),
        ("Ch31 Communication", "Ch19 Reaching Beyond Your Walls",
         "Ch20 Roles and Responsibilities"),
        ("Ch32 SRE Engagement", "Ch18, Ch21 Organizational Change", "\u2014"),
    ]),
    "Part V \u2014 \uacb0\ub860 \ubc0f \ubd80\ub85d": ("V", [
        ("Ch33 Lessons Learned", "\u2014", "Ch21 Building a Culture"),
        ("Ch34 Conclusion", "Conclusion", "Conclusion"),
        ("Appendix A-F", "Appendix A-C", "Appendix A"),
    ]),
}


def main() -> None:
    """메인 실행."""
    logging.basicConfig(level=logging.INFO, format="%(message)s")

    logger.info("=" * 60)
    logger.info("SRE Study Guide PPTX Builder - Glassmorphism")
    logger.info("=" * 60)

    prs = Presentation()
    prs.slide_width = GlassTheme.SLIDE_WIDTH
    prs.slide_height = GlassTheme.SLIDE_HEIGHT

    builder = GlassSlideBuilder(prs)

    # 1. 커버
    builder.add_cover()
    logger.info("  [1/6] \ucee4\ubc84 \uc2ac\ub77c\uc774\ub4dc")

    # 2. 개요 (3 슬라이드)
    builder.add_overview()
    builder.add_method()
    builder.add_curriculum_overview()
    logger.info("  [2/6] \uac1c\uc694 \uc2ac\ub77c\uc774\ub4dc (3\uc7a5)")

    # 3. 주차별 진도표 (12 슬라이드)
    for week_num, title, rows in WEEKS:
        builder.add_week_slide(week_num, title, rows)
    logger.info("  [3/6] \uc8fc\ucc28\ubcc4 \uc9c4\ub3c4\ud45c (12\uc7a5)")

    # 4. 교차 참조 테이블 (6 슬라이드)
    for part_name, (part_num, rows) in CROSS_REFS.items():
        builder.add_cross_ref_slide(part_name, part_num, rows)
    logger.info("  [4/6] \uad50\ucc28 \ucc38\uc870 \ud14c\uc774\ube14 (6\uc7a5)")

    # 5. 보충 학습 (2 슬라이드)
    builder.add_supplement_slide("\ubcf4\ucda9 \ud559\uc2b5 \u2014 BSRS \uace0\uc720 \ucc55\ud130", [
        ("Ch02 Understanding Adversaries", "\uc704\ud611 \ud589\uc704\uc790 \uc720\ud615\uacfc \ub3d9\uae30 \ubd84\uc11d"),
        ("Ch03 Case Study: Safe Proxies", "\uc548\uc804\ud55c \ud504\ub85d\uc2dc\ub97c \ud1b5\ud55c \uc811\uadfc \uc81c\uc5b4"),
        ("Ch05 Design for Least Privilege", "\ucd5c\uc18c \uad8c\ud55c \uc6d0\uce59 \uc801\uc6a9"),
        ("Ch11 Publicly Trusted CA", "\uacf5\uac1c \uc778\uc99d \uae30\uad00 \uc124\uacc4/\uc6b4\uc601"),
        ("Ch16 Disaster Planning", "\uc7ac\ud574 \ub300\ube44 \uacc4\ud68d \uc218\ub9bd"),
        ("Ch19 Chrome Security Team", "Chrome \ubcf4\uc548\ud300 \uc6b4\uc601 \uc0ac\ub840"),
    ])
    builder.add_supplement_slide("\ubcf4\ucda9 \ud559\uc2b5 \u2014 Workbook \uace0\uc720 \ucc55\ud130", [
        ("Ch12 Non-Abstract Large System Design", "\ube44\ucd94\uc0c1\uc801 \ub300\uaddc\ubaa8 \uc2dc\uc2a4\ud15c \uc124\uacc4 (NALSD)"),
        ("Ch14 Configuration Design", "\uc124\uc815 \uad00\ub9ac \uc124\uacc4 \uc6d0\uce59"),
        ("Ch15 Configuration Specifics", "\uc124\uc815 \uad00\ub9ac \uad6c\uccb4\uc801 \ud328\ud134"),
    ])
    logger.info("  [5/6] \ubcf4\ucda9 \ud559\uc2b5 (2\uc7a5)")

    # 6. 학습 경로 + 클로징
    builder.add_study_paths()
    builder.add_closing()
    logger.info("  [6/6] \ud559\uc2b5 \uacbd\ub85c + \ud074\ub85c\uc9d5 (2\uc7a5)")

    # 저장
    OUTPUT_DIR.mkdir(exist_ok=True)
    output_path = OUTPUT_DIR / "SRE_Study_Guide_Glassmorphism.pptx"
    prs.save(str(output_path))

    slide_count = len(prs.slides)
    size_kb = output_path.stat().st_size / 1024

    logger.info("")
    logger.info("=" * 60)
    logger.info("\uc0dd\uc131 \uc644\ub8cc: %s", output_path.name)
    logger.info("  \uc2ac\ub77c\uc774\ub4dc: %d\uc7a5", slide_count)
    logger.info("  \ud30c\uc77c \ud06c\uae30: %.1f KB", size_kb)
    logger.info("  \ucd9c\ub825 \uacbd\ub85c: %s", output_path)
    logger.info("=" * 60)


if __name__ == "__main__":
    main()

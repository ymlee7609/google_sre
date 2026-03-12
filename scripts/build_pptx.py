#!/usr/bin/env python3
"""SRE Book PPTX 학습 교재 빌더 - Nordic Minimalism 디자인.

docs/sre-book/summaries/ 의 한국어 요약 파일을
Part별 5권의 PPTX 학습 교재로 변환합니다.
"""

import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path

from lxml.etree import SubElement
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# --- 경로 상수 ---
BASE_DIR = Path(__file__).resolve().parent.parent
DOCS_DIR = BASE_DIR / "docs" / "sre-book"
SUMMARIES_DIR = DOCS_DIR / "summaries"
METADATA_PATH = DOCS_DIR / "_metadata.json"
OUTPUT_DIR = BASE_DIR / "output"


# --- Nordic Minimalism 디자인 상수 ---
class NordicTheme:
    """Nordic Minimalism 디자인 사양."""

    # 색상
    BG_CREAM = RGBColor(0xF4, 0xF1, 0xEC)
    BLOB_GREY = RGBColor(0xD9, 0xCF, 0xC4)
    TEXT_DARK = RGBColor(0x3D, 0x35, 0x30)
    TEXT_TAUPE = RGBColor(0x8A, 0x7A, 0x6A)
    DIVIDER = RGBColor(0xD9, 0xCF, 0xC4)

    # 폰트
    FONT_TITLE = "NanumMyeongjo"
    FONT_BODY = "Malgun Gothic"
    FONT_CAPTION = "Consolas"

    # 슬라이드 크기 (16:9 표준)
    SLIDE_WIDTH = Emu(12192000)   # 13.333"
    SLIDE_HEIGHT = Emu(6858000)   # 7.5"

    # 여백
    MARGIN_LEFT = Inches(1.2)
    CONTENT_WIDTH = Inches(10.933)

    # Blob (우하단 타원)
    BLOB_WIDTH = Inches(5.0)
    BLOB_HEIGHT = Inches(4.0)
    BLOB_TRANSPARENCY = 60  # 퍼센트

    # 3-dot accent (좌상단)
    DOT_SIZE = Inches(0.12)
    DOT_SPACING = Inches(0.22)
    DOT_LEFT = Inches(0.5)
    DOT_TOP = Inches(0.4)

    # 하단 수평선
    LINE_TOP = Inches(6.8)
    LINE_LEFT = Inches(0.5)
    LINE_WIDTH = Inches(12.333)
    LINE_HEIGHT = Pt(0.75)

    # 하단 캡션
    CAPTION_TOP = Inches(6.9)
    CAPTION_LEFT = Inches(0.5)
    CAPTION_WIDTH = Inches(10.0)


# --- 데이터 모델 ---
@dataclass
class ChapterData:
    """챕터 데이터."""

    number: int
    title: str
    chapter_type: str
    part_label: str
    principles: list[str] = field(default_factory=list)
    practices: list[str] = field(default_factory=list)
    quotes: list[str] = field(default_factory=list)
    keywords: list[str] = field(default_factory=list)

    @property
    def display_label(self) -> str:
        """슬라이드에 표시할 챕터 라벨."""
        if self.chapter_type == "appendix":
            letter = self.part_label.replace("Appendix ", "")
            return f"Appendix {letter}"
        return f"Chapter {self.number:02d}"


@dataclass
class PartInfo:
    """Part 정보."""

    number: int
    title_ko: str
    title_en: str
    chapters: list[ChapterData] = field(default_factory=list)

    @property
    def filename(self) -> str:
        """출력 파일명."""
        return f"SRE_Study_Part{self.number}_{self.title_en.replace(' ', '_')}.pptx"


# --- Markdown 파서 ---
class MarkdownParser:
    """summary.md 파일을 파싱하여 섹션별 콘텐츠 추출."""

    SECTION_MAP = {
        "핵심 원칙": "principles",
        "주요 프랙티스 및 권고사항": "practices",
        "핵심 인용": "quotes",
        "관련 키워드": "keywords",
    }

    @classmethod
    def parse_summary(cls, filepath: Path) -> dict[str, list[str]]:
        """요약 파일에서 4개 섹션의 콘텐츠 추출."""
        text = filepath.read_text(encoding="utf-8")
        sections: dict[str, list[str]] = {
            "principles": [],
            "practices": [],
            "quotes": [],
            "keywords": [],
        }

        current_section: str | None = None

        for line in text.split("\n"):
            line = line.strip()

            if line.startswith("## "):
                header = line[3:].strip()
                current_section = cls.SECTION_MAP.get(header)
                continue

            if not current_section or not line:
                continue

            if current_section == "quotes":
                if line.startswith("> "):
                    quote = line[2:].strip()
                    # 양쪽 따옴표 제거
                    quote = quote.strip('"\u201c\u201d')
                    if quote:
                        sections["quotes"].append(quote)
            elif current_section == "keywords":
                keywords = [k.strip() for k in line.split(",") if k.strip()]
                sections["keywords"].extend(keywords)
            elif line.startswith("- "):
                item = line[2:].strip()
                if item:
                    sections[current_section].append(item)

        return sections

    @staticmethod
    def get_summary_path(chapter_meta: dict) -> Path:
        """챕터 메타데이터에서 summary 파일 경로 생성."""
        ch_file = chapter_meta["file"]
        ch_type = chapter_meta.get("type", "chapter")

        if ch_type == "appendix":
            letter = ch_file.split("_")[1]  # 'a', 'b', ...
            return SUMMARIES_DIR / f"appendix_{letter}_summary.md"

        ch_num = ch_file.split("_")[0]  # 'ch01', 'ch02', ...
        return SUMMARIES_DIR / f"{ch_num}_summary.md"


# --- 슬라이드 빌더 ---
class SlideBuilder:
    """개별 슬라이드를 Nordic Minimalism 스타일로 생성."""

    def __init__(self, prs: Presentation) -> None:
        self.prs = prs
        self.blank_layout = prs.slide_layouts[6]

    # --- 공통 디자인 요소 ---

    @staticmethod
    def _set_fill_transparency(shape, transparency_pct: int) -> None:
        """도형 채우기에 투명도 설정 (0=불투명, 100=완전투명)."""
        # OOXML alpha: 100000=완전불투명, 0=완전투명
        alpha_val = str((100 - transparency_pct) * 1000)
        for solid_fill in shape._element.iter(qn("a:solidFill")):
            for color_el in solid_fill:
                alpha_el = SubElement(color_el, qn("a:alpha"))
                alpha_el.set("val", alpha_val)
            break

    def _add_background(self, slide) -> None:
        """Warm cream 배경색 적용."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = NordicTheme.BG_CREAM

    def _add_blob(self, slide) -> None:
        """우하단 organic blob (타원, 60% 투명)."""
        left = Inches(9.833)   # 우측으로 1.5" 넘김
        top = Inches(4.5)     # 하단으로 1.0" 넘김

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top,
            NordicTheme.BLOB_WIDTH, NordicTheme.BLOB_HEIGHT,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NordicTheme.BLOB_GREY
        shape.line.fill.background()
        self._set_fill_transparency(shape, NordicTheme.BLOB_TRANSPARENCY)

    def _add_dots(self, slide) -> None:
        """좌상단 3-dot accent."""
        for i in range(3):
            left = NordicTheme.DOT_LEFT + NordicTheme.DOT_SPACING * i
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, NordicTheme.DOT_TOP,
                NordicTheme.DOT_SIZE, NordicTheme.DOT_SIZE,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = NordicTheme.TEXT_DARK
            shape.line.fill.background()

    def _add_bottom_line(self, slide) -> None:
        """하단 수평선."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            NordicTheme.LINE_LEFT, NordicTheme.LINE_TOP,
            NordicTheme.LINE_WIDTH, NordicTheme.LINE_HEIGHT,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NordicTheme.DIVIDER
        shape.line.fill.background()

    def _add_caption(self, slide, part_number: int, part_title: str) -> None:
        """하단 캡션."""
        txbox = slide.shapes.add_textbox(
            NordicTheme.CAPTION_LEFT, NordicTheme.CAPTION_TOP,
            NordicTheme.CAPTION_WIDTH, Inches(0.4),
        )
        p = txbox.text_frame.paragraphs[0]
        p.text = f"Site Reliability Engineering  \u2014  Part {part_number}: {part_title}"
        p.font.name = NordicTheme.FONT_CAPTION
        p.font.size = Pt(9)
        p.font.color.rgb = NordicTheme.TEXT_TAUPE

    def _add_common(self, slide, part_number: int, part_title: str) -> None:
        """모든 슬라이드 공통 요소 적용."""
        self._add_background(slide)
        self._add_blob(slide)
        self._add_dots(slide)
        self._add_bottom_line(slide)
        self._add_caption(slide, part_number, part_title)

    def _add_section_divider(self, slide, y: float) -> None:
        """얇은 구분선 추가."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            NordicTheme.MARGIN_LEFT, Inches(y),
            Inches(2.0), Pt(0.75),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NordicTheme.DIVIDER
        shape.line.fill.background()

    @staticmethod
    def _dynamic_font_size(item_count: int) -> Pt:
        """항목 수에 따른 동적 폰트 크기."""
        if item_count <= 5:
            return Pt(14)
        if item_count <= 8:
            return Pt(13)
        return Pt(12)

    # --- 슬라이드 유형별 생성 ---

    def add_cover_slide(self, part: PartInfo) -> None:
        """커버 슬라이드."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_common(slide, part.number, part.title_ko)

        # "Site Reliability Engineering"
        tb1 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(2.0),
            NordicTheme.CONTENT_WIDTH, Inches(1.0),
        )
        p1 = tb1.text_frame.paragraphs[0]
        p1.text = "Site Reliability Engineering"
        p1.font.name = NordicTheme.FONT_TITLE
        p1.font.size = Pt(42)
        p1.font.color.rgb = NordicTheme.TEXT_DARK

        # "학습 자료"
        tb2 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(3.2),
            NordicTheme.CONTENT_WIDTH, Inches(0.6),
        )
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = "\ud559\uc2b5 \uc790\ub8cc"
        p2.font.name = NordicTheme.FONT_BODY
        p2.font.size = Pt(18)
        p2.font.color.rgb = NordicTheme.TEXT_TAUPE

        # Part 제목
        tb3 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(4.2),
            NordicTheme.CONTENT_WIDTH, Inches(1.0),
        )
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = f"Part {part.number}: {part.title_ko}"
        p3.font.name = NordicTheme.FONT_TITLE
        p3.font.size = Pt(36)
        p3.font.color.rgb = NordicTheme.TEXT_DARK

    def add_overview_slide(
        self, part: PartInfo, chapters: list[ChapterData],
        page: int = 0, total_pages: int = 1,
    ) -> None:
        """Part 개요(목차) 슬라이드."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_common(slide, part.number, part.title_ko)

        # 제목
        title_text = f"Part {part.number}: {part.title_ko}"
        if total_pages > 1:
            title_text += f" ({page + 1}/{total_pages})"

        tb = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(0.8),
            NordicTheme.CONTENT_WIDTH, Inches(0.8),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.name = NordicTheme.FONT_TITLE
        p.font.size = Pt(28)
        p.font.color.rgb = NordicTheme.TEXT_DARK

        self._add_section_divider(slide, 1.6)

        # 챕터 목록
        tb2 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(1.9),
            Inches(9.0), Inches(4.5),
        )
        tf = tb2.text_frame
        tf.word_wrap = True
        font_size = Pt(14) if len(chapters) <= 10 else Pt(12)

        for i, ch in enumerate(chapters):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"{ch.display_label}  \u2014  {ch.title}"
            run.font.name = NordicTheme.FONT_BODY
            run.font.size = font_size
            run.font.color.rgb = NordicTheme.TEXT_DARK

    def add_chapter_title_slide(self, ch: ChapterData, part: PartInfo) -> None:
        """챕터 제목 슬라이드 (번호 + 영문제목 + 핵심 원칙 상위 3개)."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_common(slide, part.number, part.title_ko)

        # Chapter 번호
        tb1 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(1.5),
            NordicTheme.CONTENT_WIDTH, Inches(0.5),
        )
        p1 = tb1.text_frame.paragraphs[0]
        p1.text = ch.display_label
        p1.font.name = NordicTheme.FONT_CAPTION
        p1.font.size = Pt(14)
        p1.font.color.rgb = NordicTheme.TEXT_TAUPE

        # 영문 제목
        tb2 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(2.1),
            Inches(9.0), Inches(1.2),
        )
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = ch.title
        p2.font.name = NordicTheme.FONT_TITLE
        p2.font.size = Pt(36)
        p2.font.color.rgb = NordicTheme.TEXT_DARK

        self._add_section_divider(slide, 3.5)

        # 핵심 원칙 상위 3개
        top_principles = ch.principles[:3]
        if top_principles:
            tb3 = slide.shapes.add_textbox(
                NordicTheme.MARGIN_LEFT, Inches(3.9),
                Inches(9.0), Inches(2.5),
            )
            tf = tb3.text_frame
            tf.word_wrap = True
            for i, principle in enumerate(top_principles):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                run = p.add_run()
                run.text = f"\u2014  {principle}"
                run.font.name = NordicTheme.FONT_BODY
                run.font.size = Pt(13)
                run.font.color.rgb = NordicTheme.TEXT_TAUPE

    def add_principles_slide(self, ch: ChapterData, part: PartInfo) -> None:
        """핵심 원칙 슬라이드."""
        if not ch.principles:
            return

        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_common(slide, part.number, part.title_ko)

        # 헤더 (Chapter 라벨 + 섹션 제목)
        tb = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(0.8),
            NordicTheme.CONTENT_WIDTH, Inches(0.6),
        )
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{ch.display_label}  |  "
        r1.font.name = NordicTheme.FONT_CAPTION
        r1.font.size = Pt(11)
        r1.font.color.rgb = NordicTheme.TEXT_TAUPE
        r2 = p.add_run()
        r2.text = "\ud575\uc2ec \uc6d0\uce59"
        r2.font.name = NordicTheme.FONT_TITLE
        r2.font.size = Pt(24)
        r2.font.color.rgb = NordicTheme.TEXT_DARK

        self._add_section_divider(slide, 1.5)

        # 원칙 목록
        font_size = self._dynamic_font_size(len(ch.principles))
        tb2 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(1.8),
            Inches(9.5), Inches(4.8),
        )
        tf = tb2.text_frame
        tf.word_wrap = True

        for i, principle in enumerate(ch.principles):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"\u2014  {principle}"
            run.font.name = NordicTheme.FONT_BODY
            run.font.size = font_size
            run.font.color.rgb = NordicTheme.TEXT_DARK

    def add_practices_slide(
        self, ch: ChapterData, part: PartInfo,
        items: list[str], page: int = 0, total_pages: int = 1,
    ) -> None:
        """프랙티스 슬라이드."""
        if not items:
            return

        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_common(slide, part.number, part.title_ko)

        # 헤더
        tb = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(0.8),
            NordicTheme.CONTENT_WIDTH, Inches(0.6),
        )
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{ch.display_label}  |  "
        r1.font.name = NordicTheme.FONT_CAPTION
        r1.font.size = Pt(11)
        r1.font.color.rgb = NordicTheme.TEXT_TAUPE

        title = "\uc8fc\uc694 \ud504\ub799\ud2f0\uc2a4 \ubc0f \uad8c\uace0\uc0ac\ud56d"
        if total_pages > 1:
            title += f" ({page + 1}/{total_pages})"
        r2 = p.add_run()
        r2.text = title
        r2.font.name = NordicTheme.FONT_TITLE
        r2.font.size = Pt(24)
        r2.font.color.rgb = NordicTheme.TEXT_DARK

        self._add_section_divider(slide, 1.5)

        # 프랙티스 목록
        font_size = self._dynamic_font_size(len(items))
        tb2 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(1.8),
            Inches(9.5), Inches(4.8),
        )
        tf = tb2.text_frame
        tf.word_wrap = True

        for i, practice in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"\u2014  {practice}"
            run.font.name = NordicTheme.FONT_BODY
            run.font.size = font_size
            run.font.color.rgb = NordicTheme.TEXT_DARK

    def add_quote_keywords_slide(self, ch: ChapterData, part: PartInfo) -> None:
        """인용/키워드 슬라이드."""
        if not ch.quotes and not ch.keywords:
            return

        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_common(slide, part.number, part.title_ko)

        y_cursor = 1.0  # 인치 단위

        # 인용문
        if ch.quotes:
            # 큰따옴표 장식
            tb_qm = slide.shapes.add_textbox(
                NordicTheme.MARGIN_LEFT, Inches(y_cursor),
                Inches(1.0), Inches(0.8),
            )
            p_qm = tb_qm.text_frame.paragraphs[0]
            p_qm.text = "\u201c"
            p_qm.font.name = NordicTheme.FONT_TITLE
            p_qm.font.size = Pt(48)
            p_qm.font.color.rgb = NordicTheme.DIVIDER

            y_cursor += 0.6

            for quote in ch.quotes[:3]:
                tb_q = slide.shapes.add_textbox(
                    Inches(1.8), Inches(y_cursor),
                    Inches(8.5), Inches(1.0),
                )
                tb_q.text_frame.word_wrap = True
                p_q = tb_q.text_frame.paragraphs[0]
                p_q.text = quote
                p_q.font.name = NordicTheme.FONT_TITLE
                p_q.font.size = Pt(15)
                p_q.font.italic = True
                p_q.font.color.rgb = NordicTheme.TEXT_DARK
                y_cursor += 1.0

        # 키워드
        if ch.keywords:
            if ch.quotes:
                self._add_section_divider(slide, y_cursor)
                y_cursor += 0.4

            keywords_text = "  |  ".join(ch.keywords)
            tb_kw = slide.shapes.add_textbox(
                NordicTheme.MARGIN_LEFT, Inches(y_cursor),
                Inches(10.0), Inches(1.0),
            )
            tb_kw.text_frame.word_wrap = True
            p_kw = tb_kw.text_frame.paragraphs[0]
            p_kw.text = keywords_text
            p_kw.font.name = NordicTheme.FONT_CAPTION
            p_kw.font.size = Pt(10)
            p_kw.font.color.rgb = NordicTheme.TEXT_TAUPE

    def add_closing_slide(self, part: PartInfo) -> None:
        """클로징 슬라이드."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_common(slide, part.number, part.title_ko)

        tb1 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(2.8),
            NordicTheme.CONTENT_WIDTH, Inches(1.0),
        )
        p1 = tb1.text_frame.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = f"End of Part {part.number}"
        p1.font.name = NordicTheme.FONT_TITLE
        p1.font.size = Pt(36)
        p1.font.color.rgb = NordicTheme.TEXT_DARK

        tb2 = slide.shapes.add_textbox(
            NordicTheme.MARGIN_LEFT, Inches(4.0),
            NordicTheme.CONTENT_WIDTH, Inches(0.6),
        )
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.text = part.title_ko
        p2.font.name = NordicTheme.FONT_BODY
        p2.font.size = Pt(18)
        p2.font.color.rgb = NordicTheme.TEXT_TAUPE


# --- 프레젠테이션 빌더 ---
class PresentationBuilder:
    """Part별 PPTX 조립."""

    # (번호, 한국어제목, 영문제목, 메타데이터 part 라벨 목록)
    PART_DEFS: list[tuple[int, str, str, list[str]]] = [
        (1, "\uc11c\ub860", "Introduction", ["I - Introduction"]),
        (2, "\uc6d0\uce59", "Principles", ["II - Principles"]),
        (3, "\uc2e4\ubb34", "Practices", ["III - Practices"]),
        (4, "\uad00\ub9ac", "Management", ["IV - Management"]),
        (5, "\uacb0\ub860 \ubc0f \ubd80\ub85d", "Conclusions_and_Appendices", [
            "V - Conclusions",
            "Appendix A", "Appendix B", "Appendix C",
            "Appendix D", "Appendix E", "Appendix F",
        ]),
    ]

    def __init__(self) -> None:
        with open(METADATA_PATH, encoding="utf-8") as f:
            self.metadata: dict = json.load(f)
        self.parser = MarkdownParser()

    def _load_chapter(self, ch_meta: dict) -> ChapterData:
        """챕터 메타데이터와 summary 파싱 결과 결합."""
        summary_path = self.parser.get_summary_path(ch_meta)

        sections: dict[str, list[str]] = {}
        if summary_path.exists():
            sections = self.parser.parse_summary(summary_path)
        else:
            logger.warning("Summary not found: %s", summary_path)

        return ChapterData(
            number=ch_meta.get("number", 0),
            title=ch_meta["title"],
            chapter_type=ch_meta.get("type", "chapter"),
            part_label=ch_meta.get("part", ""),
            principles=sections.get("principles", []),
            practices=sections.get("practices", []),
            quotes=sections.get("quotes", []),
            keywords=sections.get("keywords", []),
        )

    def build_part(self, part_def: tuple[int, str, str, list[str]]) -> Path:
        """Part별 PPTX 생성."""
        number, title_ko, title_en, part_labels = part_def
        part = PartInfo(number=number, title_ko=title_ko, title_en=title_en)

        # 해당 Part의 챕터 수집
        for ch_meta in self.metadata["chapters"]:
            if ch_meta.get("part", "") in part_labels:
                part.chapters.append(self._load_chapter(ch_meta))

        logger.info(
            "Part %d: %s (%d chapters)",
            number, title_ko, len(part.chapters),
        )

        # Presentation 생성
        prs = Presentation()
        prs.slide_width = NordicTheme.SLIDE_WIDTH
        prs.slide_height = NordicTheme.SLIDE_HEIGHT

        builder = SlideBuilder(prs)

        # 1. 커버
        builder.add_cover_slide(part)

        # 2. 개요(목차) - 챕터 12개 초과시 2페이지 분할
        chapters = part.chapters
        if len(chapters) > 12:
            mid = len(chapters) // 2
            builder.add_overview_slide(part, chapters[:mid], page=0, total_pages=2)
            builder.add_overview_slide(part, chapters[mid:], page=1, total_pages=2)
        else:
            builder.add_overview_slide(part, chapters)

        # 3. 챕터별 슬라이드
        for ch in chapters:
            builder.add_chapter_title_slide(ch, part)
            builder.add_principles_slide(ch, part)

            # 프랙티스 10개 초과시 2장 분할
            if ch.practices:
                if len(ch.practices) > 10:
                    mid = (len(ch.practices) + 1) // 2
                    builder.add_practices_slide(
                        ch, part, ch.practices[:mid], page=0, total_pages=2,
                    )
                    builder.add_practices_slide(
                        ch, part, ch.practices[mid:], page=1, total_pages=2,
                    )
                else:
                    builder.add_practices_slide(ch, part, ch.practices)

            builder.add_quote_keywords_slide(ch, part)

        # 4. 클로징
        builder.add_closing_slide(part)

        # 저장
        output_path = OUTPUT_DIR / part.filename
        prs.save(str(output_path))

        slide_count = len(prs.slides)
        logger.info("  -> %s (%d slides)", output_path.name, slide_count)
        return output_path

    def build_all(self) -> list[Path]:
        """전체 5권 빌드."""
        OUTPUT_DIR.mkdir(exist_ok=True)
        return [self.build_part(pd) for pd in self.PART_DEFS]


def main() -> None:
    """메인 실행."""
    logging.basicConfig(level=logging.INFO, format="%(message)s")

    logger.info("=" * 60)
    logger.info("SRE Book PPTX Builder - Nordic Minimalism")
    logger.info("=" * 60)

    builder = PresentationBuilder()
    results = builder.build_all()

    logger.info("")
    logger.info("=" * 60)
    logger.info("\uc0dd\uc131 \uc644\ub8cc: %d files", len(results))
    for path in results:
        size_kb = path.stat().st_size / 1024
        logger.info("  %s (%.1f KB)", path.name, size_kb)
    logger.info("\ucd9c\ub825 \ub514\ub809\ud1a0\ub9ac: %s", OUTPUT_DIR)
    logger.info("=" * 60)


if __name__ == "__main__":
    main()
